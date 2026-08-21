from converters.base import BaseConverter
import PyPDF2


class PptxConverter(BaseConverter):
    name = 'PowerPoint Presentation'
    output_format = 'pptx'

    def convert(self, input_path: str, output_path: str, options: dict = None) -> bool:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt

            presentation = Presentation()
            blank_layout = presentation.slide_layouts[6]

            with open(input_path, 'rb') as pdf_file:
                reader = PyPDF2.PdfReader(pdf_file)

                for page_number, page in enumerate(reader.pages, start=1):
                    slide = presentation.slides.add_slide(blank_layout)
                    title_box = slide.shapes.add_textbox(
                        Inches(0.6), Inches(0.4), Inches(12), Inches(0.6)
                    )
                    title_frame = title_box.text_frame
                    title_frame.text = f'Page {page_number}'
                    title_frame.paragraphs[0].font.size = Pt(24)
                    title_frame.paragraphs[0].font.bold = True

                    content_box = slide.shapes.add_textbox(
                        Inches(0.6), Inches(1.2), Inches(12), Inches(5.8)
                    )
                    content_frame = content_box.text_frame
                    content_frame.word_wrap = True
                    content_frame.text = (page.extract_text() or '').strip()
                    for paragraph in content_frame.paragraphs:
                        paragraph.font.size = Pt(16)

            presentation.save(output_path)
            return True
        except Exception as error:
            print(f'Error converting to PPTX: {error}')
            return False